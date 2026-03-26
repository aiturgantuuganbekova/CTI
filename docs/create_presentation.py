from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color palette
BG_DARK = RGBColor(0x0F, 0x0F, 0x1A)
BG_CARD = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xFF)
ACCENT_CYAN = RGBColor(0x00, 0xD4, 0xAA)
ACCENT_PURPLE = RGBColor(0x7B, 0x61, 0xFF)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY = RGBColor(0x99, 0x99, 0x99)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)

def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_gradient_bar(slide, left, top, width, height, color1, color2):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color1
    shape.line.fill.background()
    return shape

def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multiline_text(slide, left, top, width, height, lines, default_size=16, default_color=LIGHT_GRAY, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if isinstance(line_data, dict):
            p.text = line_data.get("text", "")
            p.font.size = Pt(line_data.get("size", default_size))
            p.font.color.rgb = line_data.get("color", default_color)
            p.font.bold = line_data.get("bold", False)
            p.font.name = line_data.get("font", font_name)
            p.alignment = line_data.get("align", PP_ALIGN.LEFT)
            if "space_before" in line_data:
                p.space_before = Pt(line_data["space_before"])
            if "space_after" in line_data:
                p.space_after = Pt(line_data["space_after"])
        else:
            p.text = str(line_data)
            p.font.size = Pt(default_size)
            p.font.color.rgb = default_color
            p.font.name = font_name
    return txBox

def add_slide_number(slide, num, total=15):
    add_text_box(slide, Inches(12.2), Inches(7.0), Inches(1), Inches(0.4),
                 f"{num} / {total}", font_size=10, color=MID_GRAY, alignment=PP_ALIGN.RIGHT)

def add_top_bar(slide):
    add_gradient_bar(slide, Inches(0), Inches(0), Inches(13.333), Pt(4), ACCENT_BLUE, ACCENT_PURPLE)


# ============================================================
# SLIDE 1 – Cover
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)
add_gradient_bar(slide, Inches(0), Inches(0), Inches(13.333), Pt(5), ACCENT_BLUE, ACCENT_CYAN)
add_gradient_bar(slide, Inches(0), Inches(7.5) - Pt(5), Inches(13.333), Pt(5), ACCENT_PURPLE, ACCENT_BLUE)

# Decorative shapes
circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(0.5), Inches(2.5), Inches(2.5))
circle1.fill.solid()
circle1.fill.fore_color.rgb = RGBColor(0x00, 0x96, 0xFF)
circle1.fill.fore_color.brightness = 0.85
circle1.line.fill.background()

circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.3), Inches(5.5), Inches(1.8), Inches(1.8))
circle2.fill.solid()
circle2.fill.fore_color.rgb = RGBColor(0x7B, 0x61, 0xFF)
circle2.fill.fore_color.brightness = 0.85
circle2.line.fill.background()

# Title
add_text_box(slide, Inches(1.5), Inches(1.2), Inches(10), Inches(0.6),
             "DIPLOMA PROJECT", font_size=14, color=ACCENT_CYAN, bold=True)
add_accent_line(slide, Inches(1.5), Inches(1.85), Inches(2), ACCENT_CYAN)

add_text_box(slide, Inches(1.5), Inches(2.1), Inches(10), Inches(1.2),
             "Development of a Cryptocurrency\nTrading Indicator", font_size=36, color=WHITE, bold=True)

# Info cards
info_items = [
    ("Student", "Aiturgan Tuuganbekova"),
    ("Group", "COM 22"),
    ("Program", "Computer Science"),
    ("Supervisor", "Dr. Daniiar Satybaldiev"),
]
for i, (label, value) in enumerate(info_items):
    x = Inches(1.5) + Inches(2.7) * i
    card = add_shape_bg(slide, x, Inches(4.2), Inches(2.4), Inches(1.4), BG_CARD)
    add_text_box(slide, x + Inches(0.2), Inches(4.35), Inches(2), Inches(0.4),
                 label, font_size=11, color=ACCENT_CYAN, bold=True)
    add_text_box(slide, x + Inches(0.2), Inches(4.7), Inches(2), Inches(0.7),
                 value, font_size=14, color=WHITE, bold=False)

add_text_box(slide, Inches(1.5), Inches(6.2), Inches(10), Inches(0.4),
             "2026", font_size=18, color=MID_GRAY, bold=False)


# ============================================================
# SLIDE 2 – Introduction
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide)
add_slide_number(slide, 2)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.6),
             "Introduction", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(1.5), ACCENT_BLUE)

bullets = [
    ("Rapid growth of cryptocurrency markets", "01"),
    ("High volatility and risk", "02"),
    ("Difficulty of making data-driven trading decisions", "03"),
    ("Need for automated and analytical tools", "04"),
]
for i, (text, num) in enumerate(bullets):
    y = Inches(1.6) + Inches(0.85) * i
    card = add_shape_bg(slide, Inches(0.8), y, Inches(6.5), Inches(0.7), BG_CARD)
    add_text_box(slide, Inches(1.0), y + Inches(0.12), Inches(0.5), Inches(0.45),
                 num, font_size=14, color=ACCENT_BLUE, bold=True)
    add_text_box(slide, Inches(1.6), y + Inches(0.12), Inches(5.5), Inches(0.45),
                 text, font_size=16, color=LIGHT_GRAY)

# Problem statement box
problem_box = add_shape_bg(slide, Inches(0.8), Inches(5.2), Inches(11.5), Inches(1.5), RGBColor(0x14, 0x20, 0x3A))
add_text_box(slide, Inches(1.1), Inches(5.35), Inches(3), Inches(0.4),
             "PROBLEM STATEMENT", font_size=12, color=ACCENT_ORANGE, bold=True)
add_text_box(slide, Inches(1.1), Inches(5.75), Inches(11), Inches(0.8),
             "Retail traders lack reliable, adaptive, and research-based indicators\nfor spot cryptocurrency trading.",
             font_size=18, color=WHITE, bold=False)

# Right side decorative chart mockup
chart_card = add_shape_bg(slide, Inches(8), Inches(1.2), Inches(4.5), Inches(3.6), BG_CARD)
add_text_box(slide, Inches(8.3), Inches(1.4), Inches(4), Inches(0.4),
             "Crypto Market Volatility", font_size=13, color=MID_GRAY, bold=True)
# Simple bars to represent chart
bar_heights = [1.0, 1.8, 1.2, 2.5, 0.8, 2.0, 1.5, 2.8, 1.3, 1.7]
bar_colors = [ACCENT_BLUE, ACCENT_CYAN, ACCENT_BLUE, ACCENT_CYAN, ACCENT_BLUE,
              ACCENT_CYAN, ACCENT_BLUE, ACCENT_CYAN, ACCENT_BLUE, ACCENT_CYAN]
for j, (h, c) in enumerate(zip(bar_heights, bar_colors)):
    bx = Inches(8.4) + Inches(0.38) * j
    by = Inches(4.4) - Inches(h)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bx, by, Inches(0.25), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = c
    bar.line.fill.background()


# ============================================================
# SLIDE 3 – Problem Statement
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide)
add_slide_number(slide, 3)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
             "Problem Statement", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(1.5), ACCENT_BLUE)

# Left column: Most traders rely on
col1 = add_shape_bg(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(2.5), BG_CARD)
add_text_box(slide, Inches(1.1), Inches(1.75), Inches(5), Inches(0.4),
             "Most traders rely on:", font_size=16, color=ACCENT_CYAN, bold=True)
items1 = ["Simple indicators (RSI, MACD, Moving Averages)", "Subjective chart analysis"]
for i, item in enumerate(items1):
    add_text_box(slide, Inches(1.3), Inches(2.3) + Inches(0.45) * i, Inches(4.8), Inches(0.4),
                 f"\u2022  {item}", font_size=14, color=LIGHT_GRAY)

# Right column: Many indicators
col2 = add_shape_bg(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(2.5), BG_CARD)
add_text_box(slide, Inches(7.3), Inches(1.75), Inches(5), Inches(0.4),
             "Many indicators:", font_size=16, color=ACCENT_ORANGE, bold=True)
items2 = ["Lag behind price", "Produce false signals in volatile markets",
          "Are not optimized for crypto-specific behavior"]
for i, item in enumerate(items2):
    add_text_box(slide, Inches(7.5), Inches(2.3) + Inches(0.45) * i, Inches(4.8), Inches(0.4),
                 f"\u2022  {item}", font_size=14, color=LIGHT_GRAY)

# Key problem box
kp_box = add_shape_bg(slide, Inches(0.8), Inches(4.6), Inches(11.5), Inches(2.2), RGBColor(0x1A, 0x0A, 0x0A))
add_accent_line(slide, Inches(0.8), Inches(4.6), Inches(11.5), ACCENT_ORANGE)
add_text_box(slide, Inches(1.2), Inches(4.9), Inches(3), Inches(0.4),
             "KEY PROBLEM", font_size=14, color=ACCENT_ORANGE, bold=True)
add_text_box(slide, Inches(1.2), Inches(5.4), Inches(10.8), Inches(1.2),
             "There is no adaptive, research-driven indicator specifically optimized for\nBinance Spot trading with integrated backtesting and signal validation.",
             font_size=18, color=WHITE)


# ============================================================
# SLIDE 4 – Research Objectives
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide)
add_slide_number(slide, 4)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
             "Research Objectives", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(1.5), ACCENT_BLUE)

# Main goal box
goal_box = add_shape_bg(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(1.0), RGBColor(0x0A, 0x1A, 0x2A))
add_accent_line(slide, Inches(0.8), Inches(1.5), Inches(11.5), ACCENT_CYAN)
add_text_box(slide, Inches(1.1), Inches(1.65), Inches(1.5), Inches(0.35),
             "MAIN GOAL", font_size=11, color=ACCENT_CYAN, bold=True)
add_text_box(slide, Inches(1.1), Inches(1.95), Inches(11), Inches(0.45),
             "Develop and validate a cryptocurrency trading indicator with backtesting support.",
             font_size=18, color=WHITE, bold=True)

# Objectives grid
objectives = [
    ("01", "Analyze existing technical indicators", ACCENT_BLUE),
    ("02", "Compare effectiveness in crypto markets", ACCENT_CYAN),
    ("03", "Design a new composite indicator", ACCENT_PURPLE),
    ("04", "Implement backtesting module", ACCENT_BLUE),
    ("05", "Integrate Telegram notification system", ACCENT_CYAN),
    ("06", "Evaluate performance metrics", ACCENT_PURPLE),
]
for i, (num, text, color) in enumerate(objectives):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + Inches(3.9) * col
    y = Inches(3.0) + Inches(1.5) * row
    card = add_shape_bg(slide, x, y, Inches(3.6), Inches(1.2), BG_CARD)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.15), y + Inches(0.2), Inches(0.55), Inches(0.55))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    add_text_box(slide, x + Inches(0.15), y + Inches(0.28), Inches(0.55), Inches(0.4),
                 num, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.85), y + Inches(0.25), Inches(2.6), Inches(0.7),
                 text, font_size=14, color=LIGHT_GRAY)


# ============================================================
# SLIDE 5 – Literature Review
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide)
add_slide_number(slide, 5)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
             "Literature Review", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(1.5), ACCENT_BLUE)

# Left: Widely used indicators
left_card = add_shape_bg(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.0), BG_CARD)
add_text_box(slide, Inches(1.1), Inches(1.65), Inches(5), Inches(0.4),
             "Widely Used Indicators", font_size=18, color=ACCENT_CYAN, bold=True)

indicators = [
    ("RSI", "Relative Strength Index"),
    ("MACD", "Moving Average Convergence Divergence"),
    ("BB", "Bollinger Bands"),
    ("EMA/SMA", "Exponential / Simple Moving Average"),
    ("ATR", "Average True Range (volatility)"),
]
for i, (abbr, full) in enumerate(indicators):
    y = Inches(2.2) + Inches(0.75) * i
    tag = add_shape_bg(slide, Inches(1.2), y, Inches(1.2), Inches(0.45), ACCENT_BLUE)
    add_text_box(slide, Inches(1.2), y + Inches(0.05), Inches(1.2), Inches(0.35),
                 abbr, font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2.6), y + Inches(0.05), Inches(3.5), Inches(0.35),
                 full, font_size=13, color=LIGHT_GRAY)

# Right: Research focus
right_card = add_shape_bg(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(5.0), BG_CARD)
add_text_box(slide, Inches(7.3), Inches(1.65), Inches(5), Inches(0.4),
             "Research Focus", font_size=18, color=ACCENT_PURPLE, bold=True)

focus_items = [
    "Signal lag problem",
    "Overfitting in backtesting",
    "Market regime adaptation",
    "Noise filtering techniques"
]
for i, item in enumerate(focus_items):
    y = Inches(2.3) + Inches(0.85) * i
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.4), y + Inches(0.08), Inches(0.2), Inches(0.2))
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT_PURPLE
    dot.line.fill.background()
    add_text_box(slide, Inches(7.8), y, Inches(4.5), Inches(0.4),
                 item, font_size=15, color=LIGHT_GRAY)


# ============================================================
# SLIDE 6 – Comparison of Existing Solutions
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide)
add_slide_number(slide, 6)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
             "Comparison of Existing Solutions", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(1.5), ACCENT_BLUE)

# Comparison criteria - left
crit_card = add_shape_bg(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(3.0), BG_CARD)
add_text_box(slide, Inches(1.1), Inches(1.65), Inches(5), Inches(0.4),
             "Comparison Criteria", font_size=16, color=ACCENT_BLUE, bold=True)
criteria = ["Signal accuracy", "Responsiveness", "Adaptability to volatility",
            "False positive rate", "Performance on 5m\u20134h timeframes"]
for i, c in enumerate(criteria):
    add_text_box(slide, Inches(1.3), Inches(2.15) + Inches(0.4) * i, Inches(5), Inches(0.35),
                 f"\u25B8  {c}", font_size=13, color=LIGHT_GRAY)

# Findings - right
find_card = add_shape_bg(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(3.0), BG_CARD)
add_text_box(slide, Inches(7.3), Inches(1.65), Inches(5), Inches(0.4),
             "Key Findings", font_size=16, color=ACCENT_CYAN, bold=True)
findings = [
    "RSI works better in range markets",
    "MACD works better in trending markets",
    "Single indicators fail in mixed conditions",
    "Composite approaches perform better"
]
for i, f in enumerate(findings):
    add_text_box(slide, Inches(7.5), Inches(2.15) + Inches(0.45) * i, Inches(4.8), Inches(0.4),
                 f"\u2713  {f}", font_size=13, color=LIGHT_GRAY)

# Conclusion box
conc_box = add_shape_bg(slide, Inches(0.8), Inches(5.0), Inches(11.5), Inches(1.2), RGBColor(0x0A, 0x1A, 0x2A))
add_accent_line(slide, Inches(0.8), Inches(5.0), Inches(11.5), ACCENT_CYAN)
add_text_box(slide, Inches(1.1), Inches(5.15), Inches(2), Inches(0.35),
             "CONCLUSION", font_size=11, color=ACCENT_CYAN, bold=True)
add_text_box(slide, Inches(1.1), Inches(5.5), Inches(11), Inches(0.5),
             "There is potential for hybrid or adaptive indicator design.", font_size=18, color=WHITE, bold=True)


# ============================================================
# SLIDE 7 – Proposed Solution
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide)
add_slide_number(slide, 7)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
             "Proposed Solution", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(1.5), ACCENT_BLUE)

# Main description
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.5),
             "A hybrid trading indicator for Binance Spot market that:", font_size=18, color=LIGHT_GRAY)

features = [
    ("\u26A1", "Combines trend + momentum + volatility signals", ACCENT_BLUE),
    ("\U0001F50D", "Filters noise using adaptive thresholds", ACCENT_CYAN),
    ("\u23F0", "Works on 5m\u20134h timeframes", ACCENT_PURPLE),
    ("\U0001F4CA", "Includes historical backtesting", ACCENT_BLUE),
    ("\U0001F4E9", "Sends Telegram alerts with entry price", ACCENT_CYAN),
]
for i, (icon, text, color) in enumerate(features):
    y = Inches(2.2) + Inches(0.85) * i
    card = add_shape_bg(slide, Inches(0.8), y, Inches(7), Inches(0.7), BG_CARD)
    icon_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y + Inches(0.1), Inches(0.5), Inches(0.5))
    icon_circle.fill.solid()
    icon_circle.fill.fore_color.rgb = color
    icon_circle.line.fill.background()
    add_text_box(slide, Inches(1.7), y + Inches(0.15), Inches(5.8), Inches(0.4),
                 text, font_size=15, color=LIGHT_GRAY)

# Core idea box
core_box = add_shape_bg(slide, Inches(8.5), Inches(2.2), Inches(4), Inches(4.0), RGBColor(0x14, 0x20, 0x3A))
add_accent_line(slide, Inches(8.5), Inches(2.2), Inches(4), ACCENT_ORANGE)
add_text_box(slide, Inches(8.8), Inches(2.5), Inches(3.5), Inches(0.4),
             "CORE IDEA", font_size=14, color=ACCENT_ORANGE, bold=True)
add_text_box(slide, Inches(8.8), Inches(3.1), Inches(3.5), Inches(2.5),
             "Reduce false signals by confirming trades across multiple logic layers.",
             font_size=18, color=WHITE)


# ============================================================
# SLIDE 8 – System Architecture
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide)
add_slide_number(slide, 8)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
             "System Architecture", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(1.5), ACCENT_BLUE)

# Technology stack boxes
tech_stack = [
    ("Frontend", "React", ACCENT_CYAN),
    ("Backend", "Spring Boot\nREST API", ACCENT_BLUE),
    ("Database", "PostgreSQL", ACCENT_PURPLE),
    ("External", "Binance API\nTelegram Bot", ACCENT_ORANGE),
]
for i, (label, tech, color) in enumerate(tech_stack):
    x = Inches(0.8) + Inches(3.1) * i
    card = add_shape_bg(slide, x, Inches(1.5), Inches(2.8), Inches(1.8), BG_CARD)
    top_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.5), Inches(2.8), Pt(4))
    top_line.fill.solid()
    top_line.fill.fore_color.rgb = color
    top_line.line.fill.background()
    add_text_box(slide, x + Inches(0.2), Inches(1.7), Inches(2.4), Inches(0.35),
                 label, font_size=12, color=color, bold=True)
    add_text_box(slide, x + Inches(0.2), Inches(2.1), Inches(2.4), Inches(0.8),
                 tech, font_size=18, color=WHITE, bold=True)

# Flow diagram
flow_card = add_shape_bg(slide, Inches(0.8), Inches(3.8), Inches(11.5), Inches(3.0), BG_CARD)
add_text_box(slide, Inches(1.1), Inches(3.95), Inches(3), Inches(0.35),
             "SYSTEM FLOW", font_size=12, color=ACCENT_BLUE, bold=True)

flow_items = ["User", "Frontend", "REST API", "Indicator\nEngine", "Database"]
flow_items2 = ["", "", "", "Backtesting\nEngine", "Telegram\nModule"]
for i, item in enumerate(flow_items):
    x = Inches(1.0) + Inches(2.3) * i
    box = add_shape_bg(slide, x, Inches(4.5), Inches(1.6), Inches(0.9), RGBColor(0x14, 0x20, 0x3A))
    top_l = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(4.5), Inches(1.6), Pt(3))
    top_l.fill.solid()
    top_l.fill.fore_color.rgb = ACCENT_BLUE
    top_l.line.fill.background()
    add_text_box(slide, x, Inches(4.6), Inches(1.6), Inches(0.7),
                 item, font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Arrow
    if i < len(flow_items) - 1:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(1.7), Inches(4.75), Inches(0.5), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ACCENT_CYAN
        arrow.line.fill.background()

# Second row
for i in range(3, 5):
    if flow_items2[i]:
        x = Inches(1.0) + Inches(2.3) * i
        box2 = add_shape_bg(slide, x, Inches(5.7), Inches(1.6), Inches(0.9), RGBColor(0x14, 0x20, 0x3A))
        top_l2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(5.7), Inches(1.6), Pt(3))
        top_l2.fill.solid()
        top_l2.fill.fore_color.rgb = ACCENT_PURPLE
        top_l2.line.fill.background()
        add_text_box(slide, x, Inches(5.8), Inches(1.6), Inches(0.7),
                     flow_items2[i], font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 9 – Indicator Logic
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide)
add_slide_number(slide, 9)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
             "Indicator Logic (Conceptual)", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(1.5), ACCENT_BLUE)

logic_layers = [
    ("1", "Trend Filter", "EMA-based direction detection", ACCENT_BLUE),
    ("2", "Momentum Confirmation", "RSI or custom oscillator", ACCENT_CYAN),
    ("3", "Volatility Filter", "ATR-based threshold adjustment", ACCENT_PURPLE),
    ("4", "Entry Trigger", "Condition alignment check", ACCENT_ORANGE),
    ("5", "Exit Logic", "Stop-loss and take-profit rules", RGBColor(0xFF, 0x44, 0x44)),
]

for i, (num, title, desc, color) in enumerate(logic_layers):
    y = Inches(1.6) + Inches(1.05) * i
    # Connector line
    if i < len(logic_layers) - 1:
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.45), y + Inches(0.8), Pt(3), Inches(0.3))
        line.fill.solid()
        line.fill.fore_color.rgb = MID_GRAY
        line.line.fill.background()

    card = add_shape_bg(slide, Inches(0.8), y, Inches(7.5), Inches(0.85), BG_CARD)
    # Number
    num_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y + Inches(0.15), Inches(0.55), Inches(0.55))
    num_circle.fill.solid()
    num_circle.fill.fore_color.rgb = color
    num_circle.line.fill.background()
    add_text_box(slide, Inches(1.0), y + Inches(0.2), Inches(0.55), Inches(0.45),
                 num, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.8), y + Inches(0.1), Inches(3), Inches(0.35),
                 title, font_size=16, color=WHITE, bold=True)
    add_text_box(slide, Inches(1.8), y + Inches(0.45), Inches(5.5), Inches(0.35),
                 desc, font_size=13, color=MID_GRAY)

# Signal note
note_box = add_shape_bg(slide, Inches(9), Inches(2.0), Inches(3.8), Inches(3.5), RGBColor(0x14, 0x20, 0x3A))
add_accent_line(slide, Inches(9), Inches(2.0), Inches(3.8), ACCENT_CYAN)
add_text_box(slide, Inches(9.3), Inches(2.3), Inches(3.3), Inches(0.4),
             "SIGNAL RULE", font_size=13, color=ACCENT_CYAN, bold=True)
add_text_box(slide, Inches(9.3), Inches(2.9), Inches(3.3), Inches(2.2),
             "Signal is generated ONLY when ALL conditions align.\n\nThis multi-layer approach significantly reduces false positives.",
             font_size=15, color=WHITE)


# ============================================================
# SLIDE 10 – Backtesting Module
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide)
add_slide_number(slide, 10)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
             "Backtesting Module", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(1.5), ACCENT_BLUE)

# Features
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5), Inches(0.4),
             "Backtesting Features", font_size=18, color=ACCENT_CYAN, bold=True)

bt_features = [
    "Historical candle data from Binance",
    "Parameter optimization",
    "Visualization of trades",
]
for i, feat in enumerate(bt_features):
    y = Inches(2.1) + Inches(0.5) * i
    add_text_box(slide, Inches(1.0), y, Inches(5), Inches(0.4),
                 f"\u25B8  {feat}", font_size=15, color=LIGHT_GRAY)

# Performance metrics cards
metrics = [
    ("Win Rate", "%", ACCENT_BLUE),
    ("Profit Factor", "ratio", ACCENT_CYAN),
    ("Max Drawdown", "%", ACCENT_ORANGE),
    ("Sharpe Ratio", "risk-adj.", ACCENT_PURPLE),
]
for i, (name, unit, color) in enumerate(metrics):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + Inches(3.0) * col
    y = Inches(4.0) + Inches(1.3) * row
    card = add_shape_bg(slide, x, y, Inches(2.7), Inches(1.1), BG_CARD)
    top_l = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(2.7), Pt(3))
    top_l.fill.solid()
    top_l.fill.fore_color.rgb = color
    top_l.line.fill.background()
    add_text_box(slide, x + Inches(0.2), y + Inches(0.15), Inches(2.3), Inches(0.4),
                 name, font_size=16, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.55), Inches(2.3), Inches(0.35),
                 unit, font_size=12, color=MID_GRAY)

# Purpose box
purpose_box = add_shape_bg(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(5.2), RGBColor(0x14, 0x20, 0x3A))
add_accent_line(slide, Inches(7), Inches(1.5), Inches(5.5), ACCENT_PURPLE)
add_text_box(slide, Inches(7.3), Inches(1.8), Inches(5), Inches(0.4),
             "PURPOSE", font_size=14, color=ACCENT_PURPLE, bold=True)
add_text_box(slide, Inches(7.3), Inches(2.4), Inches(5), Inches(2.0),
             "Validate strategy scientifically\nbefore live use.\n\nEnsure reliability through\ndata-driven evaluation.",
             font_size=18, color=WHITE)

# Mockup chart area
chart_area = add_shape_bg(slide, Inches(7.3), Inches(4.5), Inches(4.8), Inches(1.8), RGBColor(0x0A, 0x15, 0x25))
add_text_box(slide, Inches(7.5), Inches(4.6), Inches(4.4), Inches(0.3),
             "Backtest Result Preview", font_size=10, color=MID_GRAY)
# Mini chart bars
heights = [0.3, 0.5, 0.4, 0.8, 0.6, 1.0, 0.7, 0.9, 0.5, 0.8, 1.1, 0.6]
for j, h in enumerate(heights):
    bx = Inches(7.6) + Inches(0.37) * j
    by = Inches(6.1) - Inches(h)
    c = ACCENT_CYAN if h > 0.6 else ACCENT_BLUE
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bx, by, Inches(0.22), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = c
    bar.line.fill.background()


# ============================================================
# SLIDE 11 – MVP
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide)
add_slide_number(slide, 11)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
             "MVP (Minimum Viable Product)", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(1.5), ACCENT_BLUE)

mvp_items = [
    ("Indicator Engine", "Calculation and signal generation", ACCENT_BLUE),
    ("REST API", "Endpoints for frontend communication", ACCENT_CYAN),
    ("Binance Integration", "Spot data feed and pair support", ACCENT_PURPLE),
    ("Backtesting", "Historical validation and metrics", ACCENT_ORANGE),
]
for i, (title, desc, color) in enumerate(mvp_items):
    x = Inches(0.8) + Inches(3.1) * i
    card = add_shape_bg(slide, x, Inches(1.6), Inches(2.8), Inches(2.0), BG_CARD)
    top_l = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.6), Inches(2.8), Pt(4))
    top_l.fill.solid()
    top_l.fill.fore_color.rgb = color
    top_l.line.fill.background()
    add_text_box(slide, x + Inches(0.2), Inches(1.85), Inches(2.4), Inches(0.4),
                 title, font_size=16, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.2), Inches(2.25), Inches(2.4), Inches(0.8),
                 desc, font_size=13, color=MID_GRAY)

# Telegram alerts card
tg_card = add_shape_bg(slide, Inches(0.8), Inches(4.2), Inches(11.5), Inches(2.5), BG_CARD)
add_text_box(slide, Inches(1.1), Inches(4.35), Inches(5), Inches(0.4),
             "Telegram Alert Details", font_size=18, color=ACCENT_CYAN, bold=True)

tg_items = [
    ("Entry Price", "Exact price level for trade entry"),
    ("Timeframe", "Active timeframe (5m, 15m, 1h, 4h)"),
    ("Trading Pair", "Supported Binance Spot pairs"),
    ("Signal Direction", "BUY or SELL with confidence level"),
]
for i, (label, desc) in enumerate(tg_items):
    col = i % 2
    row = i // 2
    x = Inches(1.1) + Inches(5.5) * col
    y = Inches(4.9) + Inches(0.75) * row
    dot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y + Inches(0.05), Inches(0.15), Inches(0.15))
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT_CYAN
    dot.line.fill.background()
    add_text_box(slide, x + Inches(0.3), y - Inches(0.03), Inches(1.5), Inches(0.3),
                 label, font_size=14, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.25), Inches(4.5), Inches(0.3),
                 desc, font_size=12, color=MID_GRAY)


# ============================================================
# SLIDE 12 – Design Scheme
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide)
add_slide_number(slide, 12)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
             "Design Scheme", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(1.5), ACCENT_BLUE)

# Diagram types
diagrams = [
    ("Database Schema", "Entity relationships\nand data model", ACCENT_BLUE),
    ("UML Diagram", "Class structure\nand interfaces", ACCENT_CYAN),
    ("Sequence Diagram", "Component\ninteractions", ACCENT_PURPLE),
    ("Signal Flowchart", "Generation\nprocess flow", ACCENT_ORANGE),
]
for i, (title, desc, color) in enumerate(diagrams):
    x = Inches(0.8) + Inches(3.1) * i
    card = add_shape_bg(slide, x, Inches(1.5), Inches(2.8), Inches(1.8), BG_CARD)
    top_l = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.5), Inches(2.8), Pt(4))
    top_l.fill.solid()
    top_l.fill.fore_color.rgb = color
    top_l.line.fill.background()
    add_text_box(slide, x + Inches(0.2), Inches(1.75), Inches(2.4), Inches(0.4),
                 title, font_size=15, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.2), Inches(2.15), Inches(2.4), Inches(0.8),
                 desc, font_size=12, color=MID_GRAY)

# Example flow
flow_card = add_shape_bg(slide, Inches(0.8), Inches(3.8), Inches(11.5), Inches(2.8), BG_CARD)
add_text_box(slide, Inches(1.1), Inches(3.95), Inches(4), Inches(0.35),
             "SIGNAL GENERATION FLOW", font_size=12, color=ACCENT_BLUE, bold=True)

flow_steps = [
    ("Market\nData", ACCENT_BLUE),
    ("Indicator\nCalculation", ACCENT_CYAN),
    ("Signal\nValidation", ACCENT_PURPLE),
    ("Backtest\nEvaluation", ACCENT_ORANGE),
    ("Telegram\nNotification", RGBColor(0x00, 0xCC, 0x66)),
]
for i, (step, color) in enumerate(flow_steps):
    x = Inches(1.0) + Inches(2.2) * i
    box = add_shape_bg(slide, x, Inches(4.6), Inches(1.7), Inches(1.2), RGBColor(0x14, 0x20, 0x3A))
    top_l = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(4.6), Inches(1.7), Pt(4))
    top_l.fill.solid()
    top_l.fill.fore_color.rgb = color
    top_l.line.fill.background()
    add_text_box(slide, x, Inches(4.8), Inches(1.7), Inches(0.8),
                 step, font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    if i < len(flow_steps) - 1:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(1.8), Inches(5.0), Inches(0.35), Inches(0.25))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = MID_GRAY
        arrow.line.fill.background()


# ============================================================
# SLIDE 13 – Task Plan (Timeline)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide)
add_slide_number(slide, 13)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
             "Task Plan (Timeline)", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(1.5), ACCENT_BLUE)

phases = [
    ("Phase 1", "Research", ["Literature review", "Market analysis"], ACCENT_BLUE),
    ("Phase 2", "System Design", ["Architecture design", "Indicator logic design"], ACCENT_CYAN),
    ("Phase 3", "Implementation", ["Backend development", "Frontend dashboard", "Binance integration"], ACCENT_PURPLE),
    ("Phase 4", "Testing", ["Backtesting validation", "Performance evaluation"], ACCENT_ORANGE),
    ("Phase 5", "Finalization", ["Optimization", "Documentation", "Presentation"], RGBColor(0x00, 0xCC, 0x66)),
]

for i, (phase, title, items, color) in enumerate(phases):
    x = Inches(0.8) + Inches(2.5) * i

    # Phase card
    card = add_shape_bg(slide, x, Inches(1.5), Inches(2.2), Inches(5.2), BG_CARD)

    # Phase header
    header = add_shape_bg(slide, x, Inches(1.5), Inches(2.2), Inches(0.9), color)
    add_text_box(slide, x, Inches(1.55), Inches(2.2), Inches(0.35),
                 phase, font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(1.85), Inches(2.2), Inches(0.4),
                 title, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    for j, item in enumerate(items):
        add_text_box(slide, x + Inches(0.15), Inches(2.6) + Inches(0.5) * j, Inches(1.9), Inches(0.4),
                     f"\u2022 {item}", font_size=11, color=LIGHT_GRAY)

    # Connector arrow
    if i < len(phases) - 1:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                       x + Inches(2.25), Inches(3.5), Inches(0.2), Inches(0.2))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = MID_GRAY
        arrow.line.fill.background()


# ============================================================
# SLIDE 14 – Results & Expected Impact
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide)
add_slide_number(slide, 14)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
             "Results & Expected Impact", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(1.5), ACCENT_BLUE)

# Expected outcomes
out_card = add_shape_bg(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.0), BG_CARD)
add_text_box(slide, Inches(1.1), Inches(1.65), Inches(5), Inches(0.4),
             "Expected Outcomes", font_size=18, color=ACCENT_CYAN, bold=True)

outcomes = [
    "Improved signal quality compared to\nsingle-indicator systems",
    "Scientifically validated performance",
    "Reduced emotional trading",
    "Practical tool for spot traders",
]
for i, outcome in enumerate(outcomes):
    y = Inches(2.3) + Inches(0.7) * i
    check = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), y + Inches(0.05), Inches(0.25), Inches(0.25))
    check.fill.solid()
    check.fill.fore_color.rgb = ACCENT_CYAN
    check.line.fill.background()
    add_text_box(slide, Inches(1.6), y, Inches(4.5), Inches(0.55),
                 outcome, font_size=14, color=LIGHT_GRAY)

# Impact
imp_card = add_shape_bg(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(5.0), BG_CARD)
add_text_box(slide, Inches(7.3), Inches(1.65), Inches(5), Inches(0.4),
             "Impact", font_size=18, color=ACCENT_PURPLE, bold=True)

impacts = [
    ("Educational Value", "Contributes to academic research\nin quantitative trading", ACCENT_BLUE),
    ("Research Contribution", "Novel composite indicator approach\nfor crypto markets", ACCENT_CYAN),
    ("Startup Foundation", "Potential for commercial product\ndevelopment", ACCENT_PURPLE),
]
for i, (title, desc, color) in enumerate(impacts):
    y = Inches(2.3) + Inches(1.4) * i
    imp_inner = add_shape_bg(slide, Inches(7.3), y, Inches(4.9), Inches(1.15), RGBColor(0x14, 0x20, 0x3A))
    left_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.3), y, Pt(4), Inches(1.15))
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = color
    left_bar.line.fill.background()
    add_text_box(slide, Inches(7.6), y + Inches(0.1), Inches(4.5), Inches(0.35),
                 title, font_size=15, color=WHITE, bold=True)
    add_text_box(slide, Inches(7.6), y + Inches(0.45), Inches(4.5), Inches(0.5),
                 desc, font_size=12, color=MID_GRAY)


# ============================================================
# SLIDE 15 – References
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide)
add_slide_number(slide, 15)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
             "References", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(1.5), ACCENT_BLUE)

references = [
    "Murphy, J. J. (1999). Technical Analysis of the Financial Markets: A Comprehensive Guide to Trading Methods and Applications. New York Institute of Finance.",
    "Wilder, J. W. (1978). New Concepts in Technical Trading Systems. Trend Research.",
    "Bollinger, J. (2002). Bollinger on Bollinger Bands. McGraw-Hill.",
    "Appel, G. (2005). Technical Analysis: Power Tools for Active Investors. FT Press.",
    "Binance. (2025). Binance API Documentation. https://binance-docs.github.io/apidocs/",
    "Nakamoto, S. (2008). Bitcoin: A Peer-to-Peer Electronic Cash System. https://bitcoin.org/bitcoin.pdf",
    "Fama, E. F. (1970). Efficient Capital Markets: A Review of Theory and Empirical Work. The Journal of Finance, 25(2), 383\u2013417.",
    "Aronson, D. R. (2006). Evidence-Based Technical Analysis: Applying the Scientific Method and Statistical Inference to Trading Signals. Wiley.",
    "De Prado, M. L. (2018). Advances in Financial Machine Learning. Wiley.",
    "Chan, E. P. (2013). Algorithmic Trading: Winning Strategies and Their Rationale. Wiley.",
]

ref_card = add_shape_bg(slide, Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5), BG_CARD)

for i, ref in enumerate(references):
    y = Inches(1.55) + Inches(0.5) * i
    num_txt = f"[{i+1}]"
    add_text_box(slide, Inches(1.1), y, Inches(0.5), Inches(0.4),
                 num_txt, font_size=10, color=ACCENT_BLUE, bold=True)
    add_text_box(slide, Inches(1.6), y, Inches(10.3), Inches(0.45),
                 ref, font_size=10, color=LIGHT_GRAY)

add_text_box(slide, Inches(0.8), Inches(7.0), Inches(5), Inches(0.3),
             "APA 7th Edition Format", font_size=10, color=MID_GRAY)


# ============================================================
# Save
# ============================================================
output_path = "/diplom/aiturgan/docs/Diploma_Presentation_Aiturgan.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
